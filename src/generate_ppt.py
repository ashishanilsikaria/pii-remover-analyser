import math
import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE

base_dir = os.path.dirname(__file__)


def create_presentation(data):

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_layout_title = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout_title)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Automating File Cleansing and Analysis Leveraging AI"  # type: ignore
    subtitle.text = (  # type: ignore
        "Case Study Solution Presentation"
        "\nGroup Number: 45"
        "\nTeam Members: Ashish, Aryan, Alok, Monnisha, Anamika"
    )

    slide_layout_content = prs.slide_layouts[1]

    info_slides = [
        {
            "title": "Case Study Overview",
            "content": (
                "File Cleansing: Automatically detect and remove or mask all client-specific details and PII from various documents to ensure complete anonymization.\n\n"
                "Data Standardization: Pre-process diverse file formats (.jpeg, .png, .pdf, .xlsx, etc.) into a consistent, structured format suitable for analysis. \n\n"
                "Insight Extraction: Analyze the cleansed data to identify and extract meaningful security information, such as firewall rules, IAM policy statements, or IDS/IPS logs. The final output is presented in a standardized and readable format."
            ),
        },
        {
            "title": "Approach Taken",
            "content": (
                "PII Cleansing and Anonymization: Used Microsoft Presidio (spaCy en_core_web_lg plus custom YAML recognizers) to detect and anonymize PII in text and tables, and ImageRedactorEngine to mask text-based PII in images; engines are cached for performance.\n\n"
                "Multi-format Pre-processing: A MIME-type router standardizes inputs across PNG/JPG, XLSX, PPTX, and PDF using python-pptx (text/tables/images), PyPDF2 (text/images), pandas (read_excel), and PIL (images). All sanitized content is normalized into a consistent JSON structure.\n\n"
                "Analysis and Reporting: Google Gemini is prompted as a security consultant to return strict JSON with a heading, description, and 3–4 key findings. Results are rendered in an interactive UI with bullet-point findings and exported to a polished PowerPoint summarizing each file’s insights."
            ),
        },
    ]

    for info in info_slides:
        slide = prs.slides.add_slide(slide_layout_content)
        slide.shapes.title.text = info["title"]  # type: ignore
        content_shape = slide.placeholders[1]
        content_shape.text = info["content"]  # type: ignore
        for paragraph in content_shape.text_frame.paragraphs:  # type: ignore
            paragraph.font.size = Pt(20)

    system_design_slide = prs.slides.add_slide(prs.slide_layouts[5])
    system_design_slide.shapes.title.text = "Functional Design Diagram"  # type: ignore
    img_path = os.path.join(base_dir, "assets", "system_design.png")
    left = Inches(0.5)
    top = Inches(2.5)
    height = Inches(3)
    system_design_slide.shapes.add_picture(img_path, left, top, height=height)

    ppt_headers = ["File Name", "File Type", "File Description", "Key Findings"]
    rows_per_slide = 2
    num_data_slides = math.ceil(len(data) / rows_per_slide)

    slide_layout = prs.slide_layouts[5]
    empty_slide = prs.slide_layouts[6]
    first_table_slide = True

    for i in range(num_data_slides):

        if first_table_slide:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "File Analysis Output"  # type: ignore
            first_table_slide = False
            top = Inches(1.5)
        else:
            slide = prs.slides.add_slide(empty_slide)
            top = Inches(0.5)

        start_index = i * rows_per_slide
        end_index = start_index + rows_per_slide
        slide_data = data[start_index:end_index]

        num_rows_in_table = len(slide_data) + 1
        num_cols = 4
        left = Inches(0.25)

        width = Inches(9.5)
        height = Inches(5.5)

        table = slide.shapes.add_table(
            num_rows_in_table, num_cols, left, top, width, height
        ).table

        table.columns[0].width = Inches(1.4)
        table.columns[1].width = Inches(0.8)
        table.columns[2].width = Inches(3.6)
        table.columns[3].width = Inches(3.7)

        for col_idx, header_text in enumerate(ppt_headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        table.rows[0].height = Inches(0.7)

        for row_idx, record in enumerate(slide_data, start=1):
            cell = table.cell(row_idx, 0)
            cell.text = record[0]

            cell = table.cell(row_idx, 1)
            cell.text = record[1]

            for col_idx in [0, 1]:
                cell = table.cell(row_idx, col_idx)
                tf = cell.text_frame
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.paragraphs[0].font.size = Pt(12)

            cell = table.cell(row_idx, 2)
            tf = cell.text_frame
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            p_heading = tf.paragraphs[0]
            p_heading.text = record[2]
            p_heading.font.bold = True
            p_heading.font.size = Pt(12)

            p_desc = tf.add_paragraph()
            p_desc.text = record[3]
            p_desc.font.size = Pt(12)

            cell = table.cell(row_idx, 3)
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            tf = cell.text_frame
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            first = 0
            for finding in record[4]:
                if first == 0:
                    p_finding = tf.paragraphs[0]
                    p_finding.text = f"• {finding}"
                    p_finding.level = 0
                    p_finding.font.size = Pt(11)
                    first = 1
                else:
                    p_finding = tf.add_paragraph()
                    p_finding.text = f"• {finding}"
                    p_finding.level = 0
                    p_finding.font.size = Pt(11)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
